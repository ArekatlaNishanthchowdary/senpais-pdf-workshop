"""PDF to PowerPoint via one rasterized image per slide, not text/box
reconstruction -- slides are a visual medium, not a paragraph-shaped one, so
this is the same honesty level as the existing `pdf_to_web` tool: real pages
you can present from, not editable text. Kept out of the default install like
OCR/tables/word -- see the `extras` group in pyproject.toml.
"""

from __future__ import annotations

import io
from pathlib import Path

import pypdfium2 as pdfium

from ..pages import output_path
from ..registry import Param, register

_EMU_PER_PT = 12700


@register(
    id="pdf_to_ppt",
    label="PDF to PowerPoint",
    category="Convert",
    summary="Render every page as an image, one per slide, in a PowerPoint file.",
    params=(Param("dpi", "int", "Resolution (DPI)", default=150, minimum=72, maximum=300),),
)
def pdf_to_ppt(sources: list[Path], out_dir: Path, dpi: int = 150) -> list[Path]:
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as exc:
        raise ValueError(
            'PDF to PowerPoint needs the extras install: pip install "quire[extras]".'
        ) from exc

    src = sources[0]
    pdf = pdfium.PdfDocument(str(src))
    if len(pdf) == 0:
        raise ValueError("This document has no pages.")

    width_pt, height_pt = pdf[0].get_size()
    presentation = Presentation()
    presentation.slide_width = Emu(round(width_pt * _EMU_PER_PT))
    presentation.slide_height = Emu(round(height_pt * _EMU_PER_PT))
    blank_layout = presentation.slide_layouts[6]

    for page in pdf:
        image = page.render(scale=dpi / 72).to_pil().convert("RGB")
        buf = io.BytesIO()
        image.save(buf, "PNG")
        buf.seek(0)
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            buf, 0, 0, width=presentation.slide_width, height=presentation.slide_height
        )

    target = output_path(out_dir, src.stem, ".pptx")
    presentation.save(target)
    return [target]
