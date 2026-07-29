"""Tests for the python-pptx-backed PDF to PowerPoint operation (`extras`
group). Skips if python-pptx isn't installed, same adaptive pattern used
elsewhere for optional dependencies.
"""

from __future__ import annotations

import importlib.util
from pathlib import Path

import pytest
from pypdf import PdfWriter

from quire.core.registry import REGISTRY, load_operations

load_operations()

HAS_PPTX = importlib.util.find_spec("pptx") is not None


@pytest.fixture
def sample(tmp_path: Path) -> Path:
    writer = PdfWriter()
    for _ in range(3):
        writer.add_blank_page(width=600, height=400)
    target = tmp_path / "sample.pdf"
    with target.open("wb") as fh:
        writer.write(fh)
    return target


def test_pdf_to_ppt_reports_missing_extra(sample, tmp_path):
    if HAS_PPTX:
        pytest.skip("python-pptx installed; see test_pdf_to_ppt_one_slide_per_page")
    with pytest.raises(ValueError, match="extras"):
        REGISTRY["pdf_to_ppt"].run([sample], tmp_path / "o")


def test_pdf_to_ppt_one_slide_per_page(sample, tmp_path):
    if not HAS_PPTX:
        pytest.skip("python-pptx not installed")
    from pptx import Presentation

    out = REGISTRY["pdf_to_ppt"].run([sample], tmp_path / "o", dpi=72)[0]
    assert out.suffix == ".pptx"
    prs = Presentation(str(out))
    assert len(prs.slides) == 3
    for slide in prs.slides:
        assert len(slide.shapes) == 1

    # 600x400pt page at 4:3-ish aspect, slide size should track the page size
    assert round(prs.slide_width / prs.slide_height, 2) == round(600 / 400, 2)
